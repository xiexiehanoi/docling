#!/usr/bin/env python3
"""
통합 문서 텍스트 & 이미지 추출기
- PPT/PPTX: python-pptx 직접 추출 (API 불필요)
- DOC/DOCX: python-docx 직접 추출 (API 불필요)
- XLS/XLSX: openpyxl/xlrd 직접 추출 (API 불필요)
- HWP: LibreOffice로 DOCX/PDF 변환 후 추출
- PDF: PyMuPDF 이미지 추출 + EasyOCR 텍스트 (무료)

출력 구조:
  <output_dir>/
  ├── texts.json    # 텍스트 (슬라이드/페이지/섹션/시트별 구조화)
  ├── images.json   # 이미지 메타데이터 + ref 정보
  └── images/       # 추출된 이미지 파일들
"""

import os
import sys
import json
import io
import base64
import shutil
import tempfile
import subprocess
import time
import re
import zipfile
from datetime import datetime
from pathlib import Path

# ── 의존성 확인 ──────────────────────────────────────────────

HAS_PPTX = True
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    HAS_PPTX = False

HAS_FITZ = True
try:
    import fitz
except ImportError:
    HAS_FITZ = False

HAS_PIL = True
try:
    from PIL import Image
except ImportError:
    HAS_PIL = False

HAS_DOCX = True
try:
    from docx import Document as DocxDocument
    from docx.oxml.ns import qn
except ImportError:
    HAS_DOCX = False

HAS_OPENPYXL = True
try:
    import openpyxl
except ImportError:
    HAS_OPENPYXL = False

HAS_XLRD = True
try:
    import xlrd
except ImportError:
    HAS_XLRD = False

HAS_EASYOCR = True
try:
    import easyocr
except ImportError:
    HAS_EASYOCR = False


# ── 공통 유틸리티 ────────────────────────────────────────────

def _sanitize_for_filename(text, max_len=30):
    """텍스트를 파일명에 사용 가능한 형태로 변환"""
    clean = re.sub(r'[^\w\s]', '', text, flags=re.UNICODE)
    clean = clean.strip()
    clean = re.sub(r'\s+', ' ', clean)
    if len(clean) > max_len:
        clean = clean[:max_len].rstrip()
    return clean if clean else ""


def _rename_images_by_ref(images, image_dir, suffix_key, suffix_prefix):
    """이미지 파일을 ref 기반 설명적 이름으로 변경. {old_fname: new_fname} 매핑 반환."""
    rename_map = {}
    used_names = set()

    for img in images:
        ref = img.get("ref", "")
        sanitized = _sanitize_for_filename(ref)
        if not sanitized:
            sanitized = "이미지"

        suffix_val = img.get(suffix_key, 0)
        base_name = f"{sanitized} 이미지_{suffix_prefix}{suffix_val}"

        name = base_name
        counter = 1
        while name in used_names:
            counter += 1
            name = f"{base_name}_{counter}"
        used_names.add(name)

        old_fname = img["filename"]
        ext = os.path.splitext(old_fname)[1]
        new_fname = f"{name}{ext}"

        old_path = os.path.join(image_dir, old_fname)
        new_path = os.path.join(image_dir, new_fname)

        if os.path.exists(old_path):
            os.rename(old_path, new_path)

        rename_map[old_fname] = new_fname
        img["filename"] = new_fname
        img["path"] = f"images/{new_fname}"

    return rename_map


def _update_content_image_refs(content_list, rename_map):
    """content에서 이미지 파일명 참조를 rename_map에 따라 업데이트"""
    for item in content_list:
        if item.get("type") in ("image", "image_ref"):
            old = item.get("filename", "")
            if old in rename_map:
                item["filename"] = rename_map[old]
        if item.get("type") == "group" and "items" in item:
            _update_content_image_refs(item["items"], rename_map)


def _embed_images_in_tables(content_list, images, slide_num):
    """표 영역 안의 이미지를 해당 표의 images 필드에 매핑하고 content에서 제거"""
    table_items = [item for item in content_list
                   if item.get("type") == "table" and "position" in item]
    if not table_items:
        return

    to_remove = []

    for img in images:
        if img.get("slide_num") != slide_num:
            continue
        img_pos = img.get("position", {})
        if not img_pos:
            continue

        img_cx = img_pos.get("left", 0) + img_pos.get("width", 0) / 2
        img_cy = img_pos.get("top", 0) + img_pos.get("height", 0) / 2

        for table_item in table_items:
            t_pos = table_item["position"]
            t_top = t_pos.get("top", 0)
            t_left = t_pos.get("left", 0)
            t_w = t_pos.get("width", 0)
            t_h = t_pos.get("height", 0)

            if not (t_left <= img_cx <= t_left + t_w and
                    t_top <= img_cy <= t_top + t_h):
                continue

            # 이미지가 이 표 안에 있음
            table_data = table_item.get("table", [])
            if not table_data:
                continue

            # 컬럼 판별
            if isinstance(table_data[0], dict):
                headers = list(table_data[0].keys())
                total_rows = 1 + len(table_data)  # 헤더 + 데이터 행
            else:
                headers = [f"col{i}" for i in range(len(table_data[0]))] if table_data else []
                total_rows = len(table_data)

            col_widths = table_item.get("_col_widths", [])
            row_heights = table_item.get("_row_heights", [])

            col_name = ""
            col_idx = 0
            if headers:
                if col_widths and len(col_widths) >= len(headers):
                    # 실제 컬럼 너비 사용
                    cum = t_left
                    col_idx = len(headers) - 1
                    for ci in range(len(headers)):
                        cw = col_widths[ci] if ci < len(col_widths) else 0
                        if img_cx < cum + cw:
                            col_idx = ci
                            break
                        cum += cw
                else:
                    # 폴백: 균등 분할
                    col_count = len(headers)
                    col_w = t_w / col_count if col_count else t_w
                    col_idx = int((img_cx - t_left) / col_w)
                col_idx = min(col_idx, len(headers) - 1)
                col_name = headers[col_idx]

            # 행 판별: 실제 행 높이 사용
            if row_heights:
                cum = t_top
                row_idx = len(row_heights) - 1
                for ri, rh in enumerate(row_heights):
                    if img_cy < cum + rh:
                        row_idx = ri
                        break
                    cum += rh
            else:
                row_h = t_h / total_rows if total_rows else t_h
                row_idx = int((img_cy - t_top) / row_h)
            row_idx = min(row_idx, total_rows - 1)

            if "images" not in table_item:
                table_item["images"] = []
            table_item["images"].append({
                "filename": img["filename"],
                "column": col_name,
                "column_idx": col_idx,
                "row_idx": row_idx,
            })

            # content에서 해당 이미지 아이템 제거 마킹
            for i, ci in enumerate(content_list):
                if ci.get("type") == "image" and ci.get("filename") == img["filename"]:
                    if i not in to_remove:
                        to_remove.append(i)
                    break
            break  # 한 표에 매칭되면 다음 이미지로

    # 마킹된 이미지를 content에서 제거 (역순)
    for i in sorted(to_remove, reverse=True):
        content_list.pop(i)


def _strip_positions(content_list):
    """content에서 position, _col_widths, _row_heights 필드 재귀 제거"""
    for item in content_list:
        item.pop("position", None)
        item.pop("_col_widths", None)
        item.pop("_row_heights", None)
        if item.get("type") == "group" and "items" in item:
            _strip_positions(item["items"])


def _make_descriptive_ids(content_list, slide_num):
    """content 항목에 설명적 shape_id 부여"""
    text_counter = 0
    table_counter = 0
    group_counter = 0

    for item in content_list:
        t = item.get("type", "")
        if t == "text":
            text_counter += 1
            text = item.get("text", "")
            summary = _sanitize_for_filename(text, 20)
            if summary:
                item["shape_id"] = f"{summary}_슬라이드{slide_num}"
            else:
                item["shape_id"] = f"텍스트{text_counter}_슬라이드{slide_num}"
        elif t == "image":
            # 이미지 파일명에서 확장자 제거한 것을 shape_id로
            fname = item.get("filename", "")
            item["shape_id"] = os.path.splitext(fname)[0]
        elif t == "table":
            table_counter += 1
            table = item.get("table", [])
            if table and isinstance(table[0], dict):
                headers = list(table[0].keys())[:3]
                summary = _sanitize_for_filename(" ".join(headers), 20)
                item["shape_id"] = f"표_{summary}_슬라이드{slide_num}"
            else:
                item["shape_id"] = f"표{table_counter}_슬라이드{slide_num}"
        elif t == "group":
            group_counter += 1
            item["shape_id"] = f"그룹{group_counter}_슬라이드{slide_num}"
            if "items" in item:
                _make_descriptive_ids(item["items"], slide_num)


# ── LibreOffice 유틸리티 ─────────────────────────────────────

def find_libreoffice():
    for path in [
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]:
        if path and os.path.exists(path):
            return path
    return None


def convert_ppt_to_pptx(ppt_path):
    """PPT → PPTX 변환 (LibreOffice)"""
    soffice = find_libreoffice()
    if not soffice:
        print("ERROR: LibreOffice 필요 - brew install --cask libreoffice")
        sys.exit(1)

    print(f"PPT → PPTX 변환 중: {os.path.basename(ppt_path)}")
    tmp_dir = tempfile.mkdtemp()

    try:
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "pptx", "--outdir", tmp_dir, ppt_path],
            capture_output=True, timeout=120, text=True
        )
        # 변환된 파일명: 원본 basename에서 확장자만 .pptx로 변경
        converted_name = os.path.splitext(os.path.basename(ppt_path))[0] + ".pptx"
        tmp_out = os.path.join(tmp_dir, converted_name)
        if not os.path.exists(tmp_out):
            print(f"ERROR: PPT→PPTX 변환 실패")
            if result.stderr:
                print(result.stderr)
            if result.stdout:
                print(result.stdout)
            sys.exit(1)

        pptx_path = os.path.splitext(ppt_path)[0] + ".pptx"
        shutil.copy2(tmp_out, pptx_path)
        print(f"변환 완료: {pptx_path}")
        return pptx_path
    except subprocess.TimeoutExpired:
        print("ERROR: 변환 시간 초과 (120초)")
        sys.exit(1)
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def convert_doc_to_docx(doc_path):
    """DOC → DOCX 변환 (LibreOffice)"""
    soffice = find_libreoffice()
    if not soffice:
        print("ERROR: LibreOffice 필요 - brew install --cask libreoffice")
        sys.exit(1)

    print(f"DOC → DOCX 변환 중: {os.path.basename(doc_path)}")
    tmp_dir = tempfile.mkdtemp()

    try:
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "docx", "--outdir", tmp_dir, doc_path],
            capture_output=True, timeout=120, text=True
        )
        converted_name = os.path.splitext(os.path.basename(doc_path))[0] + ".docx"
        tmp_out = os.path.join(tmp_dir, converted_name)
        if not os.path.exists(tmp_out):
            print(f"ERROR: DOC→DOCX 변환 실패")
            if result.stderr:
                print(result.stderr)
            sys.exit(1)

        docx_path = os.path.splitext(doc_path)[0] + ".docx"
        shutil.copy2(tmp_out, docx_path)
        print(f"변환 완료: {docx_path}")
        return docx_path
    except subprocess.TimeoutExpired:
        print("ERROR: 변환 시간 초과 (120초)")
        sys.exit(1)
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def convert_xls_to_xlsx(xls_path):
    """XLS → XLSX 변환 (LibreOffice)"""
    soffice = find_libreoffice()
    if not soffice:
        print("ERROR: LibreOffice 필요 - brew install --cask libreoffice")
        sys.exit(1)

    print(f"XLS → XLSX 변환 중: {os.path.basename(xls_path)}")
    tmp_dir = tempfile.mkdtemp()

    try:
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "xlsx", "--outdir", tmp_dir, xls_path],
            capture_output=True, timeout=120, text=True
        )
        converted_name = os.path.splitext(os.path.basename(xls_path))[0] + ".xlsx"
        tmp_out = os.path.join(tmp_dir, converted_name)
        if not os.path.exists(tmp_out):
            print(f"ERROR: XLS→XLSX 변환 실패")
            if result.stderr:
                print(result.stderr)
            sys.exit(1)

        xlsx_path = os.path.splitext(xls_path)[0] + ".xlsx"
        shutil.copy2(tmp_out, xlsx_path)
        print(f"변환 완료: {xlsx_path}")
        return xlsx_path
    except subprocess.TimeoutExpired:
        print("ERROR: 변환 시간 초과 (120초)")
        sys.exit(1)
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def convert_hwp_to_docx(hwp_path):
    """HWP → DOCX 변환 (LibreOffice). 실패시 PDF 경로 반환."""
    soffice = find_libreoffice()
    if not soffice:
        print("ERROR: LibreOffice 필요 - brew install --cask libreoffice")
        sys.exit(1)

    print(f"HWP → DOCX 변환 중: {os.path.basename(hwp_path)}")
    tmp_dir = tempfile.mkdtemp()

    try:
        # DOCX 변환 시도
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "docx", "--outdir", tmp_dir, hwp_path],
            capture_output=True, timeout=120, text=True
        )
        converted_name = os.path.splitext(os.path.basename(hwp_path))[0] + ".docx"
        tmp_out = os.path.join(tmp_dir, converted_name)
        if os.path.exists(tmp_out):
            docx_path = os.path.splitext(hwp_path)[0] + ".docx"
            shutil.copy2(tmp_out, docx_path)
            print(f"변환 완료: {docx_path}")
            return docx_path, "docx"

        # DOCX 실패 → PDF 폴백
        print("DOCX 변환 실패, PDF 폴백 시도...")
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", tmp_dir, hwp_path],
            capture_output=True, timeout=120, text=True
        )
        pdf_name = os.path.splitext(os.path.basename(hwp_path))[0] + ".pdf"
        pdf_out = os.path.join(tmp_dir, pdf_name)
        if os.path.exists(pdf_out):
            pdf_path = os.path.splitext(hwp_path)[0] + ".pdf"
            shutil.copy2(pdf_out, pdf_path)
            print(f"PDF 폴백 변환 완료: {pdf_path}")
            return pdf_path, "pdf"

        print("ERROR: HWP 변환 실패 (DOCX/PDF 모두)")
        if result.stderr:
            print(result.stderr)
        sys.exit(1)
    except subprocess.TimeoutExpired:
        print("ERROR: 변환 시간 초과 (120초)")
        sys.exit(1)
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


# ── PPTX 추출 ───────────────────────────────────────────────

def extract_pptx(file_path, output_dir):
    """PPTX → texts.json + images.json + images/"""
    if not HAS_PPTX:
        print("ERROR: python-pptx 필요 - pip install python-pptx")
        sys.exit(1)

    image_dir = os.path.join(output_dir, "images")
    os.makedirs(image_dir, exist_ok=True)

    prs = Presentation(file_path)
    total = len(prs.slides)

    metadata = {
        "filename": os.path.basename(file_path),
        "file_size": os.path.getsize(file_path),
        "slide_count": total,
        "extraction_date": datetime.now().isoformat(),
    }

    slides_data = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        print(f"  슬라이드 {slide_idx}/{total} 처리 중...")
        slide_result = _extract_pptx_slide(slide, slide_idx, image_dir, output_dir)
        slides_data.append(slide_result["slide_data"])

    # texts.json (이미지를 각 슬라이드 안에 포함)
    result = {"metadata": metadata, "slides": slides_data}
    texts_path = os.path.join(output_dir, "texts.json")
    with open(texts_path, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)
    print(f"저장: {texts_path}")

    return result


def _extract_pptx_slide(slide, slide_num, image_dir, output_dir):
    """슬라이드 한 장에서 텍스트/이미지/표 추출"""
    layout_name = "Unknown"
    try:
        layout_name = slide.slide_layout.name
    except Exception:
        pass

    slide_data = {
        "slide_num": slide_num,
        "layout": layout_name,
        "content": [],
    }
    images = []

    for shape_idx, shape in enumerate(slide.shapes):
        result = _process_pptx_shape(shape, slide_num, shape_idx, image_dir, output_dir)
        if result:
            slide_data["content"].append(result["content"])
            if result.get("image"):
                images.append(result["image"])

    # 이미지 ref 보강: 이미지 근처 텍스트/표 기반 ref 생성
    _enrich_image_refs(slide_data["content"], images, slide_num)

    # 이미지 파일명을 ref 기반 설명적 이름으로 변경
    rename_map = _rename_images_by_ref(images, image_dir, "slide_num", "슬라이드")
    _update_content_image_refs(slide_data["content"], rename_map)

    # 설명적 shape_id 생성
    _make_descriptive_ids(slide_data["content"], slide_num)

    # 표 안 이미지를 해당 표의 images 필드로 매핑 (position 필요하므로 strip 전에 실행)
    _embed_images_in_tables(slide_data["content"], images, slide_num)

    # position 제거 (출력에 불필요)
    _strip_positions(slide_data["content"])
    for img in images:
        img.pop("position", None)

    return {"slide_data": slide_data, "images": images}


def _process_pptx_shape(shape, slide_num, shape_idx, image_dir, output_dir):
    """Shape 하나 처리"""
    pos = {
        "left": float(shape.left) if shape.left else 0,
        "top": float(shape.top) if shape.top else 0,
        "width": float(shape.width) if shape.width else 0,
        "height": float(shape.height) if shape.height else 0,
    }
    shape_id = f"slide{slide_num}_shape{shape_idx}"

    # 테이블
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        table_data = _extract_table(shape)
        text = shape.text.strip() if hasattr(shape, "text") else ""
        # 실제 컬럼 너비 / 행 높이 추출 (이미지 셀 매핑용, EMU 단위)
        table_obj = shape.table
        col_widths = [float(col.width) for col in table_obj.columns]
        row_heights = [float(row.height) for row in table_obj.rows]
        return {
            "content": {
                "type": "table",
                "shape_id": shape_id,
                "position": pos,
                "text": text,
                "table": table_data,
                "_col_widths": col_widths,
                "_row_heights": row_heights,
            }
        }

    # 이미지
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            blob = shape.image.blob
            content_type = shape.image.content_type
            ext_map = {
                "image/jpeg": "jpg", "image/png": "png",
                "image/gif": "gif", "image/bmp": "bmp",
                "image/webp": "webp", "image/tiff": "tiff",
            }
            ext = ext_map.get(content_type, "png")
            fname = f"slide{slide_num:03d}_shape{shape_idx}.{ext}"
            fpath = os.path.join(image_dir, fname)

            with open(fpath, "wb") as f:
                f.write(blob)

            img_meta = {
                "slide_num": slide_num,
                "filename": fname,
                "path": f"images/{fname}",
                "position": pos,
                "size_bytes": len(blob),
                "ref": shape_id,  # 기본 ref, 나중에 보강됨
            }

            return {
                "content": {
                    "type": "image",
                    "shape_id": shape_id,
                    "position": pos,
                    "filename": fname,
                },
                "image": img_meta,
            }
        except Exception as e:
            print(f"  이미지 추출 실패 (slide{slide_num}, shape{shape_idx}): {e}")
            return None

    # 그룹
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        group_items = []
        group_images = []
        for sub_idx, subshape in enumerate(shape.shapes):
            sub_result = _process_pptx_shape(
                subshape, slide_num, f"{shape_idx}g{sub_idx}", image_dir, output_dir
            )
            if sub_result:
                group_items.append(sub_result["content"])
                if sub_result.get("image"):
                    group_images.append(sub_result["image"])

        if group_items:
            content = {
                "type": "group",
                "shape_id": shape_id,
                "position": pos,
                "items": group_items,
            }
            # 그룹 내 이미지는 첫 번째 결과에 포함
            return {"content": content, "image": group_images[0] if group_images else None}
        return None

    # 텍스트
    if hasattr(shape, "text") and shape.text.strip():
        return {
            "content": {
                "type": "text",
                "shape_id": shape_id,
                "position": pos,
                "text": shape.text.strip(),
            }
        }

    return None


def _extract_table(table_shape):
    """표 데이터를 딕셔너리 배열로 추출"""
    try:
        table = table_shape.table
        rows = []
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            rows.append(row_data)

        if len(rows) >= 2:
            headers = rows[0]
            return [
                {headers[i]: row[i] for i in range(min(len(headers), len(row)))}
                for row in rows[1:]
            ]
        return rows
    except Exception as e:
        print(f"  표 추출 실패: {e}")
        return []


def _enrich_image_refs(content_list, images, slide_num):
    """이미지 ref를 주변 텍스트/표 컬럼/섹션 기반으로 보강"""
    text_items = []
    table_items = []

    for item in content_list:
        if item.get("type") == "text":
            text_items.append(item)
        elif item.get("type") == "table":
            table_items.append(item)
        elif item.get("type") == "group":
            for sub in item.get("items", []):
                if sub.get("type") == "text":
                    text_items.append(sub)
                elif sub.get("type") == "table":
                    table_items.append(sub)

    # 섹션 마커 찾기
    section_markers = ["▣", "▶", "●", "■", "◆", "□", "○", "△", "▲"]
    section_text = None
    for ti in text_items:
        t = ti.get("text", "")
        if any(t.startswith(m) for m in section_markers):
            section_text = t[:80]
            break

    for img in images:
        if img["slide_num"] != slide_num:
            continue

        img_pos = img.get("position", {})
        img_top = img_pos.get("top", 0)
        img_left = img_pos.get("left", 0)
        img_width = img_pos.get("width", 0)
        img_height = img_pos.get("height", 0)
        img_cx = img_left + img_width / 2
        img_cy = img_top + img_height / 2

        # 1순위: 이미지가 표 영역 안에 있으면 해당 컬럼 헤더 사용
        table_ref = None
        for ti in table_items:
            t_pos = ti.get("position", {})
            t_top = t_pos.get("top", 0)
            t_left = t_pos.get("left", 0)
            t_w = t_pos.get("width", 0)
            t_h = t_pos.get("height", 0)

            if (t_left <= img_cx <= t_left + t_w and
                    t_top <= img_cy <= t_top + t_h):
                table_data = ti.get("table", [])
                if table_data and isinstance(table_data[0], dict):
                    headers = list(table_data[0].keys())
                    if headers:
                        col_widths = ti.get("_col_widths", [])
                        if col_widths and len(col_widths) >= len(headers):
                            # 실제 컬럼 너비 사용
                            cum = t_left
                            col_idx = len(headers) - 1
                            for ci in range(len(headers)):
                                cw = col_widths[ci] if ci < len(col_widths) else 0
                                if img_cx < cum + cw:
                                    col_idx = ci
                                    break
                                cum += cw
                        else:
                            col_count = len(headers)
                            col_w = t_w / col_count if col_count else t_w
                            col_idx = int((img_cx - t_left) / col_w)
                        col_idx = min(col_idx, len(headers) - 1)
                        table_ref = headers[col_idx].strip()
                break

        if table_ref:
            img["ref"] = table_ref
            img["ref_type"] = "table_column_header"
            continue

        # 2순위: 가장 가까운 텍스트
        nearest_above = None
        nearest_above_dist = float("inf")
        nearest_below = None
        nearest_below_dist = float("inf")

        for ti in text_items:
            ti_top = ti.get("position", {}).get("top", 0)
            ti_text = ti.get("text", "").strip()
            if not ti_text:
                continue

            if ti_top <= img_top:
                dist = img_top - ti_top
                if dist < nearest_above_dist:
                    nearest_above_dist = dist
                    nearest_above = ti_text[:80]
            else:
                dist = ti_top - img_top
                if dist < nearest_below_dist:
                    nearest_below_dist = dist
                    nearest_below = ti_text[:80]

        if nearest_above:
            img["ref"] = nearest_above
            img["ref_type"] = "nearest_above_text"
        elif section_text:
            img["ref"] = section_text
            img["ref_type"] = "slide_section"
        elif nearest_below:
            img["ref"] = nearest_below
            img["ref_type"] = "nearest_below_text"
        else:
            img["ref"] = f"슬라이드{slide_num} 이미지"
            img["ref_type"] = "slide_position"


# ── PDF 추출 ────────────────────────────────────────────────

def extract_pdf(file_path, output_dir):
    """PDF → texts.json + images.json + images/ (EasyOCR)"""
    if not HAS_FITZ:
        print("ERROR: PyMuPDF 필요 - pip install PyMuPDF")
        sys.exit(1)

    if not HAS_EASYOCR:
        print("ERROR: EasyOCR 필요 - pip install easyocr")
        sys.exit(1)

    # EasyOCR Reader 초기화 (한글, 영어 지원)
    print("EasyOCR 초기화 중... (최초 실행 시 모델 다운로드)")
    reader = easyocr.Reader(['ko', 'en'], gpu=False)

    image_dir = os.path.join(output_dir, "images")
    os.makedirs(image_dir, exist_ok=True)

    doc = fitz.open(file_path)
    total = len(doc)

    metadata = {
        "filename": os.path.basename(file_path),
        "file_size": os.path.getsize(file_path),
        "page_count": total,
        "extraction_date": datetime.now().isoformat(),
    }

    # 이미지 추출
    print(f"이미지 추출 중... ({total}페이지)")
    all_images = []
    for i in range(total):
        page = doc[i]
        for idx, info in enumerate(page.get_images(full=True)):
            try:
                base_img = doc.extract_image(info[0])
                data, ext = base_img["image"], base_img["ext"]
                fname = f"page{i+1:03d}_img{idx+1:02d}.{ext}"
                fpath = os.path.join(image_dir, fname)
                with open(fpath, "wb") as f:
                    f.write(data)
                all_images.append({
                    "page": i + 1,
                    "filename": fname,
                    "path": f"images/{fname}",
                    "size_bytes": len(data),
                    "ref": f"page_{i+1}_figure_{idx+1}",
                    "ref_type": "page_position",
                })
            except Exception:
                continue

    # EasyOCR 텍스트 추출
    print(f"EasyOCR 텍스트 추출 중... ({total}페이지)")

    pages_data = []
    for i in range(total):
        print(f"  페이지 {i+1}/{total} 처리 중...")
        page = doc[i]
        pix = page.get_pixmap(dpi=200)
        img_bytes = pix.tobytes("png")

        # 임시 이미지 파일로 저장 (EasyOCR은 파일 경로 필요)
        temp_img_path = os.path.join(output_dir, f"_temp_page_{i+1}.png")
        with open(temp_img_path, "wb") as f:
            f.write(img_bytes)

        text = ""
        try:
            # EasyOCR로 텍스트 추출
            result = reader.readtext(temp_img_path)
            # result: [(bbox, text, confidence), ...]
            # bbox 순서대로 텍스트 결합
            text_lines = [item[1] for item in result]
            text = "\n".join(text_lines)
        except Exception as e:
            print(f"    OCR 실패 (페이지 {i+1}): {e}")
        finally:
            # 임시 파일 삭제
            if os.path.exists(temp_img_path):
                os.remove(temp_img_path)

        page_images = [im for im in all_images if im["page"] == i + 1]

        # 이미지 ref를 OCR 텍스트 기반으로 보강
        if text.strip() and page_images:
            lines = text.strip().split("\n")
            section = None
            for line in lines:
                stripped = line.strip()
                if stripped and (stripped.startswith("#") or
                    any(stripped.startswith(m) for m in ["▣", "▶", "●", "■", "◆"])):
                    section = stripped.lstrip("#").strip()[:80]
                    break
            if section:
                for im in page_images:
                    im["ref"] = section
                    im["ref_type"] = "page_section"

        content = []
        if text.strip():
            content.append({
                "type": "text",
                "text": text.strip(),
            })
        for im in page_images:
            content.append({
                "type": "image_ref",
                "filename": im["filename"],
                "ref": im["ref"],
            })

        pages_data.append({
            "page_num": i + 1,
            "content": content,
        })

    doc.close()

    # 이미지 설명적 파일명 변경
    rename_map = _rename_images_by_ref(all_images, image_dir, "page", "페이지")
    for pd in pages_data:
        _update_content_image_refs(pd["content"], rename_map)

    # texts.json (이미지 ref 인라인 포함)
    result = {"metadata": metadata, "pages": pages_data, "images": all_images}
    texts_path = os.path.join(output_dir, "texts.json")
    with open(texts_path, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)
    print(f"저장: {texts_path}")

    return result


# ── DOCX 추출 ───────────────────────────────────────────────

def extract_docx(file_path, output_dir, original_filename=None):
    """DOCX → texts.json + images.json + images/"""
    if not HAS_DOCX:
        print("ERROR: python-docx 필요 - pip install python-docx")
        sys.exit(1)

    image_dir = os.path.join(output_dir, "images")
    os.makedirs(image_dir, exist_ok=True)

    display_name = original_filename or os.path.basename(file_path)

    # ZIP에서 이미지 추출
    zip_images = _extract_docx_images_from_zip(file_path, image_dir)

    # 문서 파싱
    doc = DocxDocument(file_path)

    metadata = {
        "filename": display_name,
        "file_size": os.path.getsize(file_path),
        "paragraph_count": len(doc.paragraphs),
        "extraction_date": datetime.now().isoformat(),
    }

    sections, all_images = _process_docx_content(doc, zip_images, image_dir)

    # 이미지 ref 보강 + 설명적 파일명 변경
    _enrich_docx_image_refs(sections, all_images)
    rename_map = _rename_images_by_ref(all_images, image_dir, "section_idx", "섹션")
    for sec in sections:
        _update_content_image_refs(sec["content"], rename_map)

    # texts.json (이미지 ref 인라인 포함)
    result = {"metadata": metadata, "sections": sections, "images": all_images}
    texts_path = os.path.join(output_dir, "texts.json")
    with open(texts_path, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)
    print(f"저장: {texts_path}")

    return result


def _extract_docx_images_from_zip(file_path, image_dir):
    """DOCX ZIP에서 word/media/ 이미지 추출, {rId: filename} 매핑 반환"""
    images_by_rid = {}
    try:
        with zipfile.ZipFile(file_path, "r") as zf:
            media_files = [n for n in zf.namelist() if n.startswith("word/media/")]
            for idx, name in enumerate(media_files, 1):
                ext = os.path.splitext(name)[1].lower()
                if ext not in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp", ".emf", ".wmf"):
                    continue
                # EMF/WMF → PNG 변환 시도
                out_ext = ext.lstrip(".")
                if ext in (".emf", ".wmf"):
                    out_ext = "png"
                fname = f"doc_img{idx:03d}.{out_ext}"
                fpath = os.path.join(image_dir, fname)
                data = zf.read(name)

                if ext in (".emf", ".wmf") and HAS_PIL:
                    try:
                        img = Image.open(io.BytesIO(data))
                        img.save(fpath, "PNG")
                        data = open(fpath, "rb").read()
                    except Exception:
                        with open(fpath, "wb") as f:
                            f.write(data)
                else:
                    with open(fpath, "wb") as f:
                        f.write(data)

                # media 파일명을 키로 저장 (word/media/image1.png → image1.png)
                media_name = os.path.basename(name)
                images_by_rid[media_name] = {
                    "filename": fname,
                    "size_bytes": len(data),
                }
    except zipfile.BadZipFile:
        print("  경고: DOCX ZIP 구조 읽기 실패, 이미지 추출 건너뜀")
    return images_by_rid


def _process_docx_content(doc, zip_images, image_dir):
    """DOCX 단락/표를 섹션별로 구조화"""
    sections = []
    all_images = []
    current_section = {
        "section_idx": 1,
        "heading": "",
        "heading_level": 0,
        "content": [],
    }
    img_counter = 0

    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if tag == "p":
            # 단락 처리
            para = None
            for p in doc.paragraphs:
                if p._element is element:
                    para = p
                    break

            if para is None:
                continue

            # 헤딩 체크
            style_name = para.style.name if para.style else "Normal"
            if style_name.startswith("Heading"):
                try:
                    level = int(style_name.replace("Heading", "").strip())
                except ValueError:
                    level = 1
                # 새 섹션 시작
                if current_section["content"] or current_section["heading"]:
                    sections.append(current_section)
                current_section = {
                    "section_idx": len(sections) + 1,
                    "heading": para.text.strip(),
                    "heading_level": level,
                    "content": [],
                }
                continue

            # 인라인 이미지 체크
            blip_elems = element.findall(f".//{qn('a:blip')}")
            if blip_elems:
                for blip in blip_elems:
                    embed_rid = blip.get(qn("r:embed"))
                    if embed_rid:
                        # rId → media 파일명 매핑
                        img_info = _resolve_docx_image(doc, embed_rid, zip_images)
                        if img_info:
                            img_counter += 1
                            current_section["content"].append({
                                "type": "image_ref",
                                "filename": img_info["filename"],
                            })
                            all_images.append({
                                "section_idx": current_section["section_idx"],
                                "filename": img_info["filename"],
                                "path": f"images/{img_info['filename']}",
                                "size_bytes": img_info["size_bytes"],
                                "ref": f"section_{current_section['section_idx']}_figure",
                                "ref_type": "section_position",
                            })

            # 텍스트
            text = para.text.strip()
            if text:
                current_section["content"].append({
                    "type": "text",
                    "text": text,
                    "style": style_name,
                })

        elif tag == "tbl":
            # 표 처리
            for table in doc.tables:
                if table._element is element:
                    table_data = _extract_docx_table(table)
                    current_section["content"].append({
                        "type": "table",
                        "table": table_data,
                    })
                    break

    # 마지막 섹션 추가
    if current_section["content"] or current_section["heading"]:
        sections.append(current_section)

    # 섹션이 비어있으면 기본 섹션 하나 생성
    if not sections:
        sections.append(current_section)

    return sections, all_images


def _resolve_docx_image(doc, embed_rid, zip_images):
    """rId를 통해 이미지 파일명 매핑"""
    try:
        rel = doc.part.rels.get(embed_rid)
        if rel and hasattr(rel, "target_ref"):
            media_name = os.path.basename(rel.target_ref)
            if media_name in zip_images:
                return zip_images[media_name]
    except Exception:
        pass
    return None


def _extract_docx_table(table):
    """DOCX 표 데이터를 딕셔너리 배열로 추출"""
    try:
        rows = []
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            rows.append(row_data)

        if len(rows) >= 2:
            headers = rows[0]
            return [
                {headers[i]: row[i] for i in range(min(len(headers), len(row)))}
                for row in rows[1:]
            ]
        return rows
    except Exception as e:
        print(f"  표 추출 실패: {e}")
        return []


def _enrich_docx_image_refs(sections, all_images):
    """DOCX 이미지 ref를 주변 단락 텍스트 기반으로 보강"""
    for img in all_images:
        sec_idx = img.get("section_idx", 0)
        # 해당 섹션 찾기
        section = None
        for sec in sections:
            if sec["section_idx"] == sec_idx:
                section = sec
                break
        if not section:
            continue

        # 섹션 헤딩이 있으면 ref로 사용
        if section.get("heading"):
            img["ref"] = section["heading"][:80]
            img["ref_type"] = "section_heading"
            continue

        # 이미지 위치 찾기
        img_filename = img["filename"]
        img_idx = None
        for i, item in enumerate(section["content"]):
            if item.get("type") == "image_ref" and item.get("filename") == img_filename:
                img_idx = i
                break
        if img_idx is None:
            continue

        # 바로 위 텍스트 찾기
        nearest_above = None
        for i in range(img_idx - 1, -1, -1):
            if section["content"][i].get("type") == "text":
                nearest_above = section["content"][i]["text"][:80]
                break

        # 바로 아래 텍스트 찾기
        nearest_below = None
        for i in range(img_idx + 1, len(section["content"])):
            if section["content"][i].get("type") == "text":
                nearest_below = section["content"][i]["text"][:80]
                break

        if nearest_above:
            img["ref"] = nearest_above
            img["ref_type"] = "nearest_above_text"
        elif nearest_below:
            img["ref"] = nearest_below
            img["ref_type"] = "nearest_below_text"


# ── XLSX 추출 ───────────────────────────────────────────────

def extract_xlsx(file_path, output_dir):
    """XLSX → texts.json + images.json + images/"""
    if not HAS_OPENPYXL:
        print("ERROR: openpyxl 필요 - pip install openpyxl")
        sys.exit(1)

    image_dir = os.path.join(output_dir, "images")
    os.makedirs(image_dir, exist_ok=True)

    # ZIP에서 이미지 추출
    zip_images = _extract_xlsx_images_from_zip(file_path, image_dir)

    wb = openpyxl.load_workbook(file_path, data_only=True)

    metadata = {
        "filename": os.path.basename(file_path),
        "file_size": os.path.getsize(file_path),
        "sheet_count": len(wb.sheetnames),
        "extraction_date": datetime.now().isoformat(),
    }

    sheets_data = []
    all_images = []

    for sheet_idx, sheet_name in enumerate(wb.sheetnames, 1):
        print(f"  시트 {sheet_idx}/{len(wb.sheetnames)} 처리 중: {sheet_name}")
        ws = wb[sheet_name]
        sheet_result = _process_xlsx_sheet(ws, sheet_idx, sheet_name, image_dir, zip_images)
        sheets_data.append(sheet_result["sheet_data"])
        all_images.extend(sheet_result["images"])

    wb.close()

    # 이미지 설명적 파일명 변경
    _rename_images_by_ref(all_images, image_dir, "sheet_idx", "시트")

    # texts.json (이미지 ref 인라인 포함)
    result = {"metadata": metadata, "sheets": sheets_data, "images": all_images}
    texts_path = os.path.join(output_dir, "texts.json")
    with open(texts_path, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)
    print(f"저장: {texts_path}")

    return result


def _extract_xlsx_images_from_zip(file_path, image_dir):
    """XLSX ZIP에서 xl/media/ 이미지 추출"""
    images = []
    try:
        with zipfile.ZipFile(file_path, "r") as zf:
            media_files = [n for n in zf.namelist() if n.startswith("xl/media/")]
            for idx, name in enumerate(media_files, 1):
                ext = os.path.splitext(name)[1].lower()
                if ext not in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp", ".emf"):
                    continue
                out_ext = ext.lstrip(".")
                if ext == ".emf":
                    out_ext = "png"
                fname = f"xl_img{idx:03d}.{out_ext}"
                fpath = os.path.join(image_dir, fname)
                data = zf.read(name)

                if ext == ".emf" and HAS_PIL:
                    try:
                        img = Image.open(io.BytesIO(data))
                        img.save(fpath, "PNG")
                    except Exception:
                        with open(fpath, "wb") as f:
                            f.write(data)
                else:
                    with open(fpath, "wb") as f:
                        f.write(data)

                images.append({
                    "media_name": os.path.basename(name),
                    "filename": fname,
                    "size_bytes": len(data),
                })
    except zipfile.BadZipFile:
        print("  경고: XLSX ZIP 구조 읽기 실패, 이미지 추출 건너뜀")
    return images


def _process_xlsx_sheet(ws, sheet_idx, sheet_name, image_dir, zip_images):
    """시트 하나의 데이터 + 이미지 처리"""
    images = []

    # 이미지 anchor 처리
    if hasattr(ws, '_images'):
        for img_idx, ws_img in enumerate(ws._images):
            anchor_cell = ""
            if hasattr(ws_img, 'anchor') and hasattr(ws_img.anchor, '_from'):
                col = ws_img.anchor._from.col
                row = ws_img.anchor._from.row
                col_letter = openpyxl.utils.get_column_letter(col + 1)
                anchor_cell = f"{col_letter}{row + 1}"

            # ZIP에서 추출한 이미지 매칭
            matched = None
            if img_idx < len(zip_images):
                matched = zip_images[img_idx]

            if matched:
                images.append({
                    "sheet_name": sheet_name,
                    "sheet_idx": sheet_idx,
                    "filename": matched["filename"],
                    "path": f"images/{matched['filename']}",
                    "size_bytes": matched["size_bytes"],
                    "anchor_cell": anchor_cell,
                    "ref": f"{sheet_name}_{anchor_cell}" if anchor_cell else f"{sheet_name}_figure",
                    "ref_type": "anchor_cell" if anchor_cell else "sheet_position",
                })

    # 데이터 추출
    dimensions = ws.dimensions or ""
    rows = list(ws.iter_rows(values_only=True))
    merged = [str(m) for m in ws.merged_cells.ranges] if ws.merged_cells else []

    content = []
    if rows:
        # 첫 행을 헤더로 시도
        headers_raw = rows[0] if rows else []
        headers = [str(h) if h is not None else "" for h in headers_raw]

        # 헤더가 유효한지 확인 (비어있지 않은 셀이 과반수)
        non_empty = sum(1 for h in headers if h.strip())
        has_valid_headers = non_empty > len(headers) / 2 if headers else False

        if has_valid_headers and len(rows) >= 2:
            data_rows = []
            for row in rows[1:]:
                row_vals = [str(v) if v is not None else "" for v in row]
                row_dict = {
                    headers[i]: row_vals[i]
                    for i in range(min(len(headers), len(row_vals)))
                }
                data_rows.append(row_dict)
            content.append({
                "type": "data",
                "headers": headers,
                "rows": data_rows,
            })
        else:
            # 헤더 없이 raw 배열
            raw_rows = []
            for row in rows:
                raw_rows.append([str(v) if v is not None else "" for v in row])
            content.append({
                "type": "data",
                "headers": [],
                "rows": raw_rows,
            })

    # 이미지 anchor 근처 텍스트로 ref 보강
    for img in images:
        anchor = img.get("anchor_cell", "")
        if not anchor or not rows:
            continue
        try:
            cell_val = ws[anchor].value
            if cell_val:
                img["ref"] = str(cell_val)[:80]
                img["ref_type"] = "anchor_cell_text"
        except Exception:
            pass

    sheet_data = {
        "sheet_name": sheet_name,
        "sheet_idx": sheet_idx,
        "dimensions": dimensions,
        "content": content,
        "merged_cells": merged,
    }

    return {"sheet_data": sheet_data, "images": images}


# ── HWP 추출 ────────────────────────────────────────────────

def extract_hwp(file_path, output_dir):
    """HWP → DOCX 변환 후 extract_docx() 호출. 실패시 PDF 폴백."""
    original_filename = os.path.basename(file_path)
    converted_path, fmt = convert_hwp_to_docx(file_path)

    if fmt == "docx":
        print(f"DOCX로 변환 성공, DOCX 추출 진행...")
        return extract_docx(converted_path, output_dir, original_filename=original_filename)
    else:
        print(f"PDF 폴백으로 추출 진행...")
        return extract_pdf(converted_path, output_dir)


# ── 메인 ─────────────────────────────────────────────────────

def main():
    supported = (".ppt", ".pptx", ".pdf", ".doc", ".docx", ".hwp", ".xls", ".xlsx")

    if len(sys.argv) < 2:
        print("사용법: python3 doc_extract.py <파일경로> [출력디렉토리]")
        print()
        print(f"지원 형식: {', '.join(supported)}")
        print()
        print("출력 구조:")
        print("  <출력디렉토리>/")
        print("  ├── texts.json    # 텍스트 (JSON)")
        print("  ├── images.json   # 이미지 메타 + ref")
        print("  └── images/       # 이미지 파일들")
        sys.exit(1)

    file_path = os.path.abspath(sys.argv[1])

    if not os.path.exists(file_path):
        print(f"ERROR: 파일 없음 - {file_path}")
        sys.exit(1)

    ext = os.path.splitext(file_path)[1].lower()
    if ext not in supported:
        print(f"ERROR: 지원하지 않는 형식 - {ext}")
        print(f"지원: {', '.join(supported)}")
        sys.exit(1)

    # 출력 디렉토리
    if len(sys.argv) > 2:
        output_dir = os.path.abspath(sys.argv[2])
    else:
        base = os.path.splitext(file_path)[0]
        output_dir = f"{base}_extracted"

    os.makedirs(output_dir, exist_ok=True)

    print(f"입력: {file_path}")
    print(f"출력: {output_dir}")
    print(f"형식: {ext}")
    print("=" * 60)

    if ext == ".ppt":
        pptx_path = convert_ppt_to_pptx(file_path)
        extract_pptx(pptx_path, output_dir)
    elif ext == ".pptx":
        extract_pptx(file_path, output_dir)
    elif ext == ".pdf":
        extract_pdf(file_path, output_dir)
    elif ext == ".doc":
        docx_path = convert_doc_to_docx(file_path)
        extract_docx(docx_path, output_dir)
    elif ext == ".docx":
        extract_docx(file_path, output_dir)
    elif ext == ".hwp":
        extract_hwp(file_path, output_dir)
    elif ext == ".xls":
        xlsx_path = convert_xls_to_xlsx(file_path)
        extract_xlsx(xlsx_path, output_dir)
    elif ext == ".xlsx":
        extract_xlsx(file_path, output_dir)

    # 결과 요약
    print()
    print("=" * 60)
    print("추출 완료!")
    print(f"  결과: {os.path.join(output_dir, 'texts.json')}")
    print(f"  이미지 폴더: {os.path.join(output_dir, 'images')}")

    img_dir = os.path.join(output_dir, "images")
    if os.path.exists(img_dir):
        img_count = len([f for f in os.listdir(img_dir) if not f.startswith(".")])
        print(f"  추출 이미지: {img_count}개")
    print("=" * 60)


if __name__ == "__main__":
    main()
