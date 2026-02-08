import os
import re
import io
import time
import subprocess
import shutil
import tempfile

from google import genai
from PIL import Image

HAS_FITZ = True
try:
    import fitz  # PyMuPDF
except ImportError:
    HAS_FITZ = False

HAS_PPTX = True
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    HAS_PPTX = False


class GeminiDocumentExtractor:

    def __init__(self, api_key: str, model_name: str = "gemini-3.0-flash-preview"):
        self.client = genai.Client(api_key=api_key)
        self.model_name = model_name

    def extract(self, file_path: str, output_dir: str, progress_callback=None):
        os.makedirs(output_dir, exist_ok=True)
        text_dir = os.path.join(output_dir, "texts")
        image_dir = os.path.join(output_dir, "images")
        os.makedirs(text_dir, exist_ok=True)
        os.makedirs(image_dir, exist_ok=True)

        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf":
            texts, images = self._process_pdf(file_path, image_dir, progress_callback)
        elif ext == ".pptx":
            texts, images = self._process_pptx(file_path, image_dir, progress_callback)
        elif ext == ".ppt":
            texts, images = self._process_ppt(file_path, image_dir, progress_callback)
        else:
            raise ValueError(f"지원하지 않는 파일 형식: {ext}")

        # 페이지별 텍스트 저장
        label = "slide" if ext in (".pptx", ".ppt") else "page"
        for t in texts:
            path = os.path.join(text_dir, f"{label}_{t['page']:03d}.md")
            with open(path, "w", encoding="utf-8") as f:
                f.write(f"# {label.title()} {t['page']}\n\n{t['text']}")

        self._save_full_text(texts, text_dir)

        if progress_callback:
            progress_callback(1.0, "완료!")

        return texts, images

    # ── PDF ────────────────────────────────────────────────────

    def _process_pdf(self, file_path, image_dir, cb=None):
        if not HAS_FITZ:
            raise RuntimeError("PyMuPDF가 필요합니다: pip3 install PyMuPDF")

        doc = fitz.open(file_path)
        total = len(doc)
        texts, images = [], []

        for i in range(total):
            if cb:
                cb(i / total, f"PDF 페이지 {i+1}/{total} 처리 중...")

            page = doc[i]

            # 텍스트: 페이지 렌더링 → Gemini OCR
            pix = page.get_pixmap(dpi=200)
            img = Image.open(io.BytesIO(pix.tobytes("png")))
            text = self._gemini_ocr(img, i + 1)
            texts.append({"page": i + 1, "text": text})

            # 이미지: 임베디드 이미지 추출
            for idx, info in enumerate(page.get_images(full=True)):
                try:
                    base = doc.extract_image(info[0])
                    data, ext = base["image"], base["ext"]
                    fname = f"page{i+1:03d}_img{idx+1:02d}.{ext}"
                    fpath = os.path.join(image_dir, fname)
                    with open(fpath, "wb") as f:
                        f.write(data)
                    images.append({"page": i + 1, "filename": fname, "path": fpath})
                except Exception:
                    continue

        doc.close()
        return texts, images

    # ── PPTX ───────────────────────────────────────────────────

    def _process_pptx(self, file_path, image_dir, cb=None):
        if not HAS_PPTX:
            raise RuntimeError("python-pptx가 필요합니다: pip3 install python-pptx")

        prs = Presentation(file_path)
        total = len(prs.slides)
        images = []

        # 이미지 추출
        for snum, slide in enumerate(prs.slides, 1):
            for sidx, shape in enumerate(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        blob = shape.image.blob
                        ext = shape.image.content_type.split("/")[-1]
                        if ext == "jpeg":
                            ext = "jpg"
                        fname = f"slide{snum:03d}_img{sidx+1:02d}.{ext}"
                        fpath = os.path.join(image_dir, fname)
                        with open(fpath, "wb") as f:
                            f.write(blob)
                        images.append({"page": snum, "filename": fname, "path": fpath})
                    except Exception:
                        continue

        # 텍스트: Gemini File API
        if cb:
            cb(0.3, "Gemini에 PPTX 업로드 중...")
        texts = self._text_via_gemini_upload(file_path, total, cb)

        return texts, images

    # ── PPT (구형 포맷) ────────────────────────────────────────

    def _process_ppt(self, file_path, image_dir, cb=None):
        # LibreOffice로 PDF 변환 시도
        pdf_path = self._convert_with_libreoffice(file_path, "pdf")
        if pdf_path:
            if cb:
                cb(0.1, "PPT → PDF 변환 완료, PDF로 처리 중...")
            return self._process_pdf(pdf_path, image_dir, cb)

        # LibreOffice 없으면 변환 불가
        raise RuntimeError(
            ".ppt(구형 포맷)은 변환이 필요합니다. 아래 중 하나를 선택하세요:\n"
            "  1) 파일을 .pptx 또는 .pdf로 변환 후 다시 실행\n"
            "     (Keynote/PowerPoint에서 열기 → 다른 이름으로 저장)\n"
            "  2) LibreOffice 설치 (자동 변환 지원):\n"
            "     brew install --cask libreoffice"
        )

    # ── Gemini 공통 ────────────────────────────────────────────

    def _call_gemini(self, contents, max_retries=5):
        """Gemini API 호출 (503 등 일시 오류 시 자동 재시도)"""
        for attempt in range(max_retries):
            try:
                response = self.client.models.generate_content(
                    model=self.model_name,
                    contents=contents,
                )
                return response.text
            except Exception as e:
                err = str(e)
                if ("503" in err or "overloaded" in err.lower() or "UNAVAILABLE" in err) and attempt < max_retries - 1:
                    wait = 2 ** attempt  # 1, 2, 4, 8, 16초
                    print(f"  [재시도] 서버 과부하, {wait}초 후 재시도... ({attempt+1}/{max_retries})")
                    time.sleep(wait)
                else:
                    raise

    def _gemini_ocr(self, image: Image.Image, page_num: int) -> str:
        prompt = (
            f"이 문서 페이지(Page {page_num})의 모든 텍스트를 정확하게 추출해주세요.\n\n"
            "규칙:\n"
            "1. 원본의 구조와 레이아웃을 최대한 유지\n"
            "2. 표는 마크다운 표 형식으로 변환\n"
            "3. 제목/소제목 계층 구조 유지\n"
            "4. 이미지 위치는 [이미지]로 표시\n"
            "5. 추출한 텍스트만 출력 (부가 설명 금지)"
        )
        return self._call_gemini([image, prompt])

    def _text_via_gemini_upload(self, file_path, total_pages, cb=None):
        # 한글 경로 문제 우회: 임시 파일로 복사 후 업로드
        ext = os.path.splitext(file_path)[1]
        tmp = tempfile.NamedTemporaryFile(suffix=ext, delete=False)
        try:
            with open(file_path, "rb") as src:
                tmp.write(src.read())
            tmp.close()
            uploaded = self.client.files.upload(file=tmp.name)
        finally:
            os.unlink(tmp.name)

        # 처리 대기
        while uploaded.state.name == "PROCESSING":
            time.sleep(1)
            uploaded = self.client.files.get(name=uploaded.name)

        if uploaded.state.name == "FAILED":
            raise RuntimeError("Gemini 파일 업로드/처리 실패")

        if cb:
            cb(0.5, "Gemini로 텍스트 추출 중...")

        count_info = f"총 {total_pages}개 슬라이드입니다.\n" if total_pages else ""
        prompt = (
            f"이 문서의 각 페이지/슬라이드에서 모든 텍스트를 정확하게 추출해주세요.\n"
            f"{count_info}\n"
            "규칙:\n"
            "1. 각 페이지를 아래 구분자로 나누세요:\n"
            "   --- SLIDE 1 ---\n"
            "   (텍스트)\n"
            "   --- SLIDE 2 ---\n"
            "   (텍스트)\n"
            "2. 표는 마크다운 표 형식으로 변환\n"
            "3. 제목/소제목 계층 구조 유지\n"
            "4. 이미지 위치는 [이미지]로 표시\n"
            "5. 추출한 텍스트만 출력 (부가 설명 금지)"
        )

        full_text = self._call_gemini([uploaded, prompt])

        # 파싱
        texts = self._parse_slide_sections(full_text)

        # 정리
        try:
            self.client.files.delete(name=uploaded.name)
        except Exception:
            pass

        return texts

    # ── 유틸리티 ───────────────────────────────────────────────

    @staticmethod
    def _convert_with_libreoffice(src_path, target_fmt):
        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if not soffice:
            mac = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
            if os.path.exists(mac):
                soffice = mac
        if not soffice:
            return None

        # 한글 경로 문제 우회: 임시 디렉토리에서 변환
        ext = os.path.splitext(src_path)[1]
        tmp_dir = tempfile.mkdtemp()
        tmp_src = os.path.join(tmp_dir, f"input{ext}")
        shutil.copy2(src_path, tmp_src)

        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", target_fmt, "--outdir", tmp_dir, tmp_src],
                capture_output=True, timeout=120,
            )
        except Exception:
            shutil.rmtree(tmp_dir, ignore_errors=True)
            return None

        tmp_out = os.path.join(tmp_dir, f"input.{target_fmt}")
        if not os.path.exists(tmp_out):
            shutil.rmtree(tmp_dir, ignore_errors=True)
            return None

        # 원본 위치에 결과 복사
        final = os.path.splitext(src_path)[0] + f".{target_fmt}"
        shutil.copy2(tmp_out, final)
        shutil.rmtree(tmp_dir, ignore_errors=True)
        return final

    @staticmethod
    def _parse_slide_sections(full_text: str) -> list:
        sections = re.split(r"---\s*SLIDE\s*(\d+)\s*---", full_text)
        texts = []
        for i in range(1, len(sections), 2):
            num = int(sections[i])
            body = sections[i + 1].strip() if i + 1 < len(sections) else ""
            texts.append({"page": num, "text": body})
        if not texts:
            texts.append({"page": 1, "text": full_text})
        return texts

    @staticmethod
    def _save_full_text(texts, text_dir):
        with open(os.path.join(text_dir, "full_text.md"), "w", encoding="utf-8") as f:
            for t in texts:
                f.write(f"# Page {t['page']}\n\n{t['text']}\n\n---\n\n")
